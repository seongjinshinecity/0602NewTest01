# 데일리브루 경쟁사 분석 PPT 5장 — Berry & Cream 팔레트 (python-pptx)
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BERRY = RGBColor(0x6D, 0x2E, 0x46)
ROSE = RGBColor(0xA2, 0x67, 0x69)
CREAM = RGBColor(0xEC, 0xE2, 0xD0)
DARK = RGBColor(0x2B, 0x1B, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOLD = RGBColor(0xD9, 0xA8, 0x5A)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font='Georgia', anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.alignment = align
        p.line_spacing = line_spacing
        for run in p.runs:
            run.font.size, run.font.bold, run.font.name = Pt(size), bold, font
            run.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.word_wrap = True
    return sh

# ── 1. 표지 (다크 베리)
s = prs.slides.add_slide(BLANK); bg(s, BERRY)
box(s, 0.9, 2.1, 11.5, 1.1, '데일리브루 경쟁사 분석', 44, CREAM, bold=True)
box(s, 0.9, 3.3, 11.5, 0.6, '성수동 반경 500m — 우리는 어디서 이기는가', 20, RGBColor(0xE8,0xC9,0xCF))
box(s, 0.9, 6.4, 11.5, 0.5, 'DAILY BREW  ·  2026. 07. 08  ·  하버스쿨 5주차 [조사] 퀘스트', 12, ROSE)
rect(s, 0.9, 4.35, 3.7, 0.02, GOLD)  # 미세 포인트
box(s, 0.9, 4.6, 11, 0.9, '"매일 굽는 디저트, 데일리브루"', 18, GOLD, font='Georgia')

# ── 2. 시장 개요 (스탯 콜아웃)
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
box(s, 0.7, 0.5, 12, 0.8, '시장 개요 — 성수동 골목 상권', 32, BERRY, bold=True)
stats = [
    ('4곳', '반경 500m 직접 경쟁 카페', '빈브라더스 · 카페모노 · 스윗아워 · 데일리브루'),
    ('4,000~5,000원', '아메리카노 가격 박스권', '가격 경쟁 여지 없음 → 제품·경험 차별화 필수'),
    ('5.2배', '팔로워 격차 (12,500 vs 2,400)', '제품력 대비 인지도가 병목 — 홍보가 과제'),
]
for i, (num, label, desc) in enumerate(stats):
    x = 0.7 + i * 4.2
    rect(s, x, 1.8, 3.9, 3.6, WHITE)
    box(s, x + 0.3, 2.25, 3.3, 1.1, num, 40 if len(num) <= 5 else 26, BERRY, bold=True)
    box(s, x + 0.3, 3.45, 3.3, 0.6, label, 15, DARK, bold=True, font='Calibri')
    box(s, x + 0.3, 4.15, 3.3, 1.1, desc, 12, RGBColor(0x6b,0x5a,0x60), font='Calibri', line_spacing=1.15)
box(s, 0.7, 5.9, 12, 1.1, '핵심 질문 — 디저트 강자 스윗아워도 외주·냉동 베이스.\n"매장에서 매일 직접 굽는 집"이라는 자리는 아직 비어 있다.', 16, DARK, font='Calibri', line_spacing=1.2)

# ── 3. 경쟁사 비교표
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
box(s, 0.7, 0.45, 12, 0.8, '경쟁사 비교 — 4개 카페 한눈에', 32, BERRY, bold=True)
cols = ['항목', '빈브라더스', '카페모노', '스윗아워', '데일리브루 ★']
rows = [
    ['아메리카노', '5,000원', '4,000원', '4,800원', '4,500원'],
    ['시그니처', '핸드드립', '대용량 라떼 1L', '수제 마카롱·케이크', '매일 굽는 치즈케이크·크로플'],
    ['좌석', '60석', '18석', '32석', '24석'],
    ['팔로워', '8,200', '3,100', '12,500', '2,400'],
    ['강점', '원두 브랜드', '가성비·회전', '디저트 비주얼·SNS', '당일 제조 신선도·카공 친화'],
    ['약점', '디저트 외주', '디저트 없음', '외주·냉동 베이스', '인지도·주말 웨이팅'],
]
tbl = s.shapes.add_table(len(rows) + 1, 5, Inches(0.7), Inches(1.5), Inches(11.9), Inches(4.6)).table
tbl.columns[0].width = Inches(1.6)
for c in range(1, 5): tbl.columns[c].width = Inches(2.575)
for c, name in enumerate(cols):
    cell = tbl.cell(0, c)
    cell.text = name
    cell.fill.solid(); cell.fill.fore_color.rgb = BERRY
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs: r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(13), True, CREAM, 'Calibri'
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xF7,0xEF,0xE3) if ci < 4 else RGBColor(0xF3,0xDD,0xE2)  # 우리 열 강조
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size, r.font.name = Pt(11.5), 'Calibri'
                r.font.color.rgb = DARK
                r.font.bold = (ci == 0 or ci == 4)
box(s, 0.7, 6.35, 12, 0.7, '팔로워는 최하위지만, "당일 제조 디저트"는 4곳 중 유일 — 제품의 자리는 확보, 알리는 일이 남았다.', 15, BERRY, bold=True, font='Calibri')

# ── 4. 우리 차별화
s = prs.slides.add_slide(BLANK); bg(s, CREAM)
box(s, 0.7, 0.5, 12, 0.8, '데일리브루 차별화 3', 32, BERRY, bold=True)
diffs = [
    ('01', '당일 제조 디저트', '반경 내 유일한 "매장에서 매일 굽는" 카페.\n오픈 주방·굽는 시간 공개로 눈으로 증명한다.'),
    ('02', '카공 친화 공간', '콘센트·와이파이·1인석. 스윗아워(모임 위주),\n카페모노(체류 불가)와 다른 "오래 머무는" 경험.'),
    ('03', '시그니처 집중', '치즈케이크·크로플 두 축. 주방 협소 제약을\n"가짓수 대신 품질"이라는 브랜드로 전환.'),
]
for i, (num, title, desc) in enumerate(diffs):
    y = 1.6 + i * 1.75
    rect(s, 0.7, y, 11.9, 1.5, WHITE)
    box(s, 1.0, y + 0.28, 1.2, 1.0, num, 32, ROSE, bold=True)
    box(s, 2.3, y + 0.2, 3.6, 0.6, title, 19, BERRY, bold=True, font='Calibri')
    box(s, 6.1, y + 0.18, 6.2, 1.2, desc, 13, DARK, font='Calibri', line_spacing=1.2)

# ── 5. 추천 액션 (다크 마무리)
s = prs.slides.add_slide(BLANK); bg(s, BERRY)
box(s, 0.7, 0.5, 12, 0.8, '추천 액션', 32, CREAM, bold=True)
actions = [
    ('베이킹 타임 콘텐츠', '매일 굽는 장면을 릴스로 — 팔로워 격차(5.2배) 해소의 최저비용 수단', '이번 주 시작'),
    ('마이크로 인플루언서 협업', '성수 카페 투어 계정 대상 체험 협업 → [홍보] 인플루언서 퀘스트로 연결', '2주 내'),
    ('주말 웨이팅 관리', 'VoC 부정 1위 "대기시간" 대응 — 치즈케이크 주말 증량 + 예상 대기시간 안내', '즉시'),
]
for i, (title, desc, when) in enumerate(actions):
    y = 1.7 + i * 1.6
    rect(s, 0.7, y, 11.9, 1.35, RGBColor(0x7E,0x3C,0x55))
    box(s, 1.05, y + 0.15, 7.5, 0.55, f'{i+1}.  {title}', 19, CREAM, bold=True, font='Calibri')
    box(s, 1.05, y + 0.72, 8.6, 0.55, desc, 12.5, RGBColor(0xE8,0xC9,0xCF), font='Calibri')
    box(s, 10.2, y + 0.38, 2.1, 0.6, when, 14, GOLD, bold=True, align=PP_ALIGN.RIGHT, font='Calibri')
box(s, 0.7, 6.75, 12, 0.5, 'DAILY BREW — 다음 회의: 인플루언서 후보 리스트 리뷰', 12, ROSE, font='Calibri')

prs.save('competitor_report.pptx')
print('saved competitor_report.pptx (5 slides)')
